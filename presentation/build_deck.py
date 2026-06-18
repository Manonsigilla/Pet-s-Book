# -*- coding: utf-8 -*-
"""
Présentation de soutenance Pet's Book — RNCP 37273, Bloc 2.
Priorité : la grille de compétences. Couleurs et typographie de Pet's Book.
Notes de chaque diapo = glossaire des termes techniques (explications simples).
Le .pptx s'importe dans Canva (gratuit) pour modification.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "presentation", "assets")
DOC = os.path.join(ROOT, "Pets-book-version1", "documentation")
LOGO = os.path.join(ROOT, "client", "public", "images", "Pets_Book.png")

# --- Identité Pet's Book ----------------------------------------------------
DARK   = RGBColor(0x59, 0x3A, 0x28)
MED    = RGBColor(0x8C, 0x59, 0x3B)
BRAND  = RGBColor(0xBF, 0x8A, 0x6B)
LIGHT  = RGBColor(0xD9, 0xD6, 0xD2)
CREAM  = RGBColor(0xF6, 0xF1, 0xEC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x4F, 0x79, 0x42)

HEAD = "Georgia"      # serif (évoque Cinzel) — universel
BODY = "Segoe UI"     # sans lisible — universel
EMW, EMH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(EMW)
prs.slide_height = Inches(EMH)
BLANK = prs.slide_layouts[6]


# --- Helpers ----------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, text, size=16, color=DARK, bold=False, font=BODY, align=PP_ALIGN.LEFT,
         after=4, before=0, first=False, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(after); p.space_before = Pt(before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return p


def bullet(tf, text, size=15, color=DARK, lead=None, glyph="•", after=6, font=BODY):
    p = tf.add_paragraph(); p.space_after = Pt(after)
    if lead:
        r = p.add_run(); r.text = f"{glyph}  {lead}"
        r.font.size = Pt(size); r.font.bold = True; r.font.name = font; r.font.color.rgb = color
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.name = font; r2.font.color.rgb = color
    else:
        r = p.add_run(); r.text = f"{glyph}  {text}"
        r.font.size = Pt(size); r.font.name = font; r.font.color.rgb = color
    return p


def check(tf, done, why, size=14, after=9):
    p = tf.add_paragraph(); p.space_after = Pt(after)
    r = p.add_run(); r.text = "✓  "
    r.font.size = Pt(size); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = GREEN
    r2 = p.add_run(); r2.text = done
    r2.font.size = Pt(size); r2.font.bold = True; r2.font.name = BODY; r2.font.color.rgb = DARK
    if why:
        r3 = p.add_run(); r3.text = " — " + why
        r3.font.size = Pt(size - 0.5); r3.font.name = BODY; r3.font.color.rgb = MED
    return p


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def img_fit(s, path, l, t, bw, bh, border=DARK):
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(bw / w, bh / h)
    dw, dh = w * ratio, h * ratio
    pic = s.shapes.add_picture(path, Inches(l + (bw - dw) / 2), Inches(t + (bh - dh) / 2),
                               Inches(dw), Inches(dh))
    if border is not None:
        pic.line.color.rgb = border; pic.line.width = Pt(1.5)
    return pic


def shape_text(s, shape_type, l, t, w, h, text, fill, txt=WHITE, size=12, bold=True, line=None):
    sp = s.shapes.add_shape(shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = BODY; r.font.color.rgb = txt
    return sp


def connect(s, x1, y1, x2, y2, color=MED, w=1.5):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    return c


_pageno = {"n": 0}


def content(title, kicker=None):
    s = slide()
    rect(s, 0, 0, EMW, EMH, WHITE)
    rect(s, 0, 0, EMW, 1.15, BRAND)
    rect(s, 0, 1.15, EMW, 0.06, DARK)
    if kicker:
        para(tb(s, 0.55, 0.16, EMW - 1.0, 0.3), kicker.upper(), size=11, color=DARK, bold=True, first=True)
        para(tb(s, 0.55, 0.42, EMW - 1.0, 0.66), title, size=25, color=WHITE, bold=True, font=HEAD, first=True)
    else:
        para(tb(s, 0.55, 0, EMW - 1.0, 1.15, anchor=MSO_ANCHOR.MIDDLE), title, size=27,
             color=WHITE, bold=True, font=HEAD, first=True)
    _pageno["n"] += 1
    para(tb(s, 0.55, EMH - 0.42, EMW - 4.0, 0.32, anchor=MSO_ANCHOR.MIDDLE),
         "Pet's Book — réseau social dédié aux animaux", size=9, color=MED, first=True)
    rp = tb(s, EMW - 3.6, EMH - 0.42, 3.05, 0.32, anchor=MSO_ANCHOR.MIDDLE).paragraphs[0]
    rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run(); rr.text = f"RNCP 37273 · Bloc 2 · {_pageno['n']:02d}"
    rr.font.size = Pt(9); rr.font.name = BODY; rr.font.color.rgb = MED
    return s


def rncp_chip(s, text):
    rect(s, 0.7, 6.4, EMW - 1.4, 0.5, CREAM, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = tb(s, 1.0, 6.4, EMW - 2.0, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Compétences Bloc 2  ·  "
    r.font.size = Pt(11.5); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = BRAND
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(11.5); r2.font.name = BODY; r2.font.color.rgb = DARK


# ===========================================================================
# 1 · TITRE
# ===========================================================================
s = slide()
rect(s, 0, 0, EMW, EMH, DARK)
rect(s, 0, 0, 0.35, EMH, BRAND); rect(s, EMW - 0.35, 0, 0.35, EMH, BRAND)
if os.path.exists(LOGO):
    img_fit(s, LOGO, EMW / 2 - 1.1, 0.7, 2.2, 1.8, border=None)
t = tb(s, 1.2, 2.7, EMW - 2.4, 1.5, anchor=MSO_ANCHOR.MIDDLE)
para(t, "Pet's Book", size=54, color=WHITE, bold=True, font=HEAD, align=PP_ALIGN.CENTER, first=True, after=4)
para(t, "Réseau social dédié aux animaux  ·  refonte fullstack 2021 → 2026",
     size=18, color=BRAND, align=PP_ALIGN.CENTER, after=0)
t2 = tb(s, 1.2, 4.5, EMW - 2.4, 1.6)
para(t2, "Titre RNCP 37273 « Développeur web fullstack »", size=16, color=CREAM, align=PP_ALIGN.CENTER, first=True, after=3)
para(t2, "Bloc 2 — Développer des interfaces frontend pour un site / une application web",
     size=14, color=LIGHT, align=PP_ALIGN.CENTER, after=12)
para(t2, "Manon Sigaud", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, after=0)
notes(s, """
GLOSSAIRE (à expliquer si on me le demande) :
• Frontend : la partie visible d'un site, affichée dans le navigateur (ce que voit l'utilisateur).
• Fullstack : qui couvre le frontend (ce qui s'affiche) ET le backend (le serveur et les données).
• RNCP : Répertoire National des Certifications Professionnelles — le registre officiel des diplômes reconnus par l'État.
• Bloc de compétences : un ensemble cohérent de compétences évaluées ensemble ; ici le Bloc 2 porte sur les interfaces frontend.
""")

# ===========================================================================
# 2 · SOMMAIRE
# ===========================================================================
s = content("Sommaire")
items = [
    ("1", "Le projet à ses débuts", "Idée de départ & cahier des charges"),
    ("2", "Conception", "Personas, moodboard, charte graphique"),
    ("3", "Avant / après", "L'évolution du site en un coup d'œil"),
    ("4", "Choix techniques", "Technologies utilisées et pourquoi"),
    ("5", "Fonctionnement (UML)", "Diagramme de cas d'utilisation"),
    ("6", "Les compétences du Bloc 2", "Éco-conception, SCSS, HTML, JavaScript"),
    ("7", "Qualité & sécurité", "Les points forts du projet"),
    ("8", "Démonstration", "Place au live !"),
]
y = 1.42
for num, ttl, desc in items:
    rect(s, 0.85, y, 0.52, 0.52, BRAND, shape=MSO_SHAPE.OVAL)
    para(tb(s, 0.85, y, 0.52, 0.52, anchor=MSO_ANCHOR.MIDDLE), num, size=17, color=WHITE,
         bold=True, font=HEAD, align=PP_ALIGN.CENTER, first=True)
    tf = tb(s, 1.6, y - 0.04, EMW - 2.3, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, ttl, size=16, color=DARK, bold=True, first=True, after=1)
    para(tf, desc, size=11.5, color=MED, after=0)
    y += 0.66
notes(s, """
Déroulé de la présentation. J'enchaîne ensuite par une démonstration en direct (≈ 7 minutes).
Aucun terme technique particulier sur cette diapo.
""")

# ===========================================================================
# 3 · START-UP PROJECT + CAHIER DES CHARGES
# ===========================================================================
s = content("Le projet à ses débuts", "L'idée & le cahier des charges (2021)")
left = tb(s, 0.7, 1.5, 6.4, 5.2)
para(left, "L'idée de départ", size=17, color=DARK, bold=True, font=HEAD, first=True, after=6)
para(left, "« Répertorier tous les animaux domestiques sur une seule plateforme. »",
     size=14, color=MED, italic=True, after=8)
for b in [
    "un profil pour chaque animal",
    "sensibiliser à la protection : identification (puce) et stérilisation",
    "condition d'entrée : justifier si l'animal n'est pas identifié / stérilisé",
    "un espace pour les pros : vétérinaires, associations, refuges",
]:
    bullet(left, b, size=13.5)
right = tb(s, 7.4, 1.5, 5.2, 5.2)
para(right, "Le cahier des charges", size=17, color=DARK, bold=True, font=HEAD, first=True, after=6)
for lead, txt in [
    ("Type : ", "plateforme de réseau social"),
    ("Équipe : ", "projet de groupe (3 personnes), 1re formation"),
    ("Cibles : ", "propriétaires d'animaux + professionnels"),
    ("Périmètre : ", "création de comptes, géolocalisation, profils"),
    ("Prévu « pour plus tard » : ", "boutique et version multilingue"),
]:
    bullet(right, txt, lead=lead, size=13.5)
para(right, "→ Beaucoup de ces idées sont aujourd'hui réalisées (score de protection, "
            "Pages pros, boutique…).", size=12.5, color=GREEN, bold=True, before=8, after=0)
notes(s, """
GLOSSAIRE :
• Cahier des charges : document qui définit le besoin, les objectifs, les cibles et le périmètre d'un projet avant de le réaliser.
• Persona : personnage fictif représentant un type d'utilisateur cible (voir diapo suivante).
• Réseau social : plateforme où les utilisateurs créent un profil, publient et interagissent.
• Géolocalisation : utiliser la position géographique (ville) pour proposer du contenu proche de l'utilisateur.
""")

# ===========================================================================
# 4 · PERSONAS + MOODBOARD
# ===========================================================================
s = content("Conception — personas & moodboard", "Comprendre les utilisateurs & l'ambiance")
para(tb(s, 0.7, 1.4, 6.0, 0.4), "Personas (utilisateurs cibles)", size=15, color=DARK, bold=True, font=HEAD, first=True)
pp = os.path.join(DOC, "6. Personas Pet's book.png")
if os.path.exists(pp):
    img_fit(s, pp, 0.7, 1.85, 6.0, 4.8)
para(tb(s, 7.0, 1.4, 5.6, 0.4), "Moodboard (direction visuelle)", size=15, color=DARK, bold=True, font=HEAD, first=True)
mb = os.path.join(DOC, "7. Moodboard pet's book.png")
if os.path.exists(mb):
    img_fit(s, mb, 7.0, 1.85, 5.6, 4.8)
notes(s, """
GLOSSAIRE :
• Persona : personnage fictif (prénom, âge, besoins, freins) qui incarne un groupe d'utilisateurs.
  Ici : Alain, vétérinaire, et Julien, jeune propriétaire de chien. Cela aide à concevoir pour de vraies attentes.
• Moodboard (planche d'ambiance) : collage d'images, couleurs et styles qui fixe l'atmosphère visuelle souhaitée avant de dessiner le site.
""")

# ===========================================================================
# 5 · CHARTE GRAPHIQUE
# ===========================================================================
s = content("Conception — la charte graphique", "L'identité visuelle de Pet's Book")
left = tb(s, 0.7, 1.5, 5.6, 5.0)
para(left, "Une identité chaleureuse", size=16, color=DARK, bold=True, font=HEAD, first=True, after=6)
bullet(left, "palette caramel et marron, douce et naturelle", size=14)
bullet(left, "typographie : Cinzel (titres) + Quicksand (texte)", size=14)
bullet(left, "logo et icônes dessinés pour la marque", size=14)
para(left, "Conservée dans la refonte : même âme, exécution modernisée — "
           "ces couleurs sont celles de cette présentation.",
     size=12.5, color=GREEN, bold=True, before=10, after=0)
# Nuancier
para(tb(s, 6.7, 1.5, 5.9, 0.4), "Couleurs de la marque", size=15, color=DARK, bold=True, font=HEAD, first=True)
sw = [("Marron foncé", "#593A28", DARK), ("Marron moyen", "#8C593B", MED),
      ("Caramel", "#BF8A6B", BRAND), ("Beige clair", "#D9D6D2", LIGHT)]
sy = 2.1
for name, hexa, col in sw:
    rect(s, 6.7, sy, 1.1, 0.8, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         line=(DARK if col == LIGHT else None))
    tf = tb(s, 8.0, sy, 4.4, 0.8, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, name, size=13, color=DARK, bold=True, first=True, after=1)
    para(tf, hexa, size=11.5, color=MED, after=0)
    sy += 0.98
# Aperçu typo
tf = tb(s, 6.7, 6.05, 5.9, 0.8)
para(tf, "Cinzel — titres élégants", size=16, color=DARK, bold=True, font=HEAD, first=True, after=1)
para(tf, "Quicksand — texte courant, lisible et rond", size=12.5, color=MED, after=0)
notes(s, """
GLOSSAIRE :
• Charte graphique : ensemble des règles visuelles d'une marque (couleurs, polices, logo) pour une identité cohérente.
• Palette : la sélection de couleurs officielles du projet.
• Typographie : le choix et l'agencement des polices de caractères.
• Code couleur hexadécimal (ex. #BF8A6B) : notation qui désigne une couleur précise pour le web.
""")

# ===========================================================================
# 6 · MAQUETTES AVANT / APRÈS
# ===========================================================================
s = content("Avant / après", "L'évolution du site")
v1 = os.path.join(ASSETS, "Pet's Book — Version 1 (site d'origine).jpg")
v2 = os.path.join(ASSETS, "Pet's Book — Application (nouvelle version).jpg")
para(tb(s, 1.4, 1.35, 4.4, 0.4, anchor=MSO_ANCHOR.MIDDLE), "2021 — site d'origine (statique)",
     size=14, color=MED, bold=True, align=PP_ALIGN.CENTER, first=True)
if os.path.exists(v1):
    img_fit(s, v1, 1.4, 1.8, 4.4, 4.3)
para(tb(s, 7.5, 1.35, 4.4, 0.4, anchor=MSO_ANCHOR.MIDDLE), "2026 — application fullstack",
     size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
if os.path.exists(v2):
    img_fit(s, v2, 7.5, 1.8, 4.4, 4.3)
# Flèche centrale
shape_text(s, MSO_SHAPE.RIGHT_ARROW, 6.05, 3.6, 1.2, 0.7, "", BRAND)
para(tb(s, 0.7, 6.35, EMW - 1.4, 0.6, anchor=MSO_ANCHOR.MIDDLE),
     "Même identité visuelle — d'une vitrine figée à une application vivante alimentée par un vrai back-end.",
     size=13, color=DARK, align=PP_ALIGN.CENTER, first=True)
notes(s, """
GLOSSAIRE :
• Site statique : pages figées, sans base de données ; le contenu ne change pas selon l'utilisateur.
• Application (dynamique) : le contenu est généré à la demande à partir de données (comptes, publications…).
• Back-end : le serveur, invisible pour l'utilisateur, qui stocke les données et répond aux demandes du site.
""")

# ===========================================================================
# 7 · TECHNOLOGIES & POURQUOI
# ===========================================================================
s = content("Technologies utilisées — et pourquoi", "Des choix maîtrisés")
colL = tb(s, 0.7, 1.5, 6.0, 5.2)
para(colL, "Frontend (côté navigateur)", size=16, color=BRAND, bold=True, font=HEAD, first=True, after=5)
bullet(colL, "outil de build moderne, rechargement instantané pendant le dev.", lead="Vite : ", size=13.5)
bullet(colL, "du CSS plus puissant (variables, organisation, réutilisation).", lead="SCSS : ", size=13.5)
bullet(colL, "pour maîtriser le langage sans dépendre d'un framework.", lead="JavaScript ES6 : ", size=13.5)
colR = tb(s, 7.0, 1.5, 5.6, 5.2)
para(colR, "Backend (côté serveur)", size=16, color=BRAND, bold=True, font=HEAD, first=True, after=5)
bullet(colR, "l'environnement qui exécute le JavaScript côté serveur.", lead="Node.js : ", size=13.5)
bullet(colR, "construit l'API (les routes qui fournissent les données).", lead="Express : ", size=13.5)
bullet(colR, "base de données légère, sans serveur séparé — idéale pour le projet.", lead="SQLite : ", size=13.5)
bullet(colR, "sécurisent les mots de passe et les sessions.", lead="bcrypt + JWT : ", size=13.5)
para(tb(s, 0.7, 6.45, EMW - 1.4, 0.5, anchor=MSO_ANCHOR.MIDDLE),
     "Pas de Bootstrap ni de framework back-end non maîtrisé : des outils choisis et compris.",
     size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)
notes(s, """
GLOSSAIRE :
• Vite : outil moderne qui assemble (« build ») le site et fournit un serveur de développement rapide.
• Build : étape qui transforme le code source en fichiers optimisés prêts pour le navigateur.
• Framework : « boîte à outils » de code prête à l'emploi (ex. React, Bootstrap). J'ai choisi de ne pas en utiliser pour le JS et le CSS afin de démontrer ma maîtrise.
• API : interface par laquelle le site demande/reçoit des données du serveur (souvent au format JSON).
• Node.js : permet d'exécuter du JavaScript en dehors du navigateur, côté serveur.
• Express : bibliothèque Node qui simplifie la création d'un serveur et de ses « routes » (URL).
• Base de données : système qui stocke et organise les données (utilisateurs, animaux…).
• SQLite : base de données rangée dans un simple fichier, sans serveur dédié.
• bcrypt : fonction qui « hache » les mots de passe (les transforme en empreinte illisible).
• JWT (jeton) : petit jeton signé prouvant qu'un utilisateur est connecté.
""")

# ===========================================================================
# 8 · DIAGRAMME UML (cas d'utilisation)
# ===========================================================================
s = content("Comment fonctionne le site — diagramme UML", "Cas d'utilisation : qui peut faire quoi")
# Cadre système
rect(s, 2.8, 1.5, 7.7, 5.1, CREAM, line=BRAND, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
para(tb(s, 3.0, 1.6, 4.0, 0.35), "Système — Pet's Book", size=12, color=MED, bold=True, first=True)
# Acteurs
a_vis = shape_text(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 2.15, 1.85, 0.62, "👤  Visiteur", BRAND, size=13)
a_mem = shape_text(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 4.0, 1.85, 0.62, "👤  Membre", MED, size=13)
a_adm = shape_text(s, MSO_SHAPE.ROUNDED_RECTANGLE, 10.75, 4.9, 1.95, 0.62, "👤  Administrateur", DARK, size=12.5)
# Cas d'utilisation (colonne centrale)
ucs = [
    "S'inscrire / se connecter",
    "Gérer ses profils animaux",
    "Publier & se faire des copains",
    "S'inscrire à un évènement & laisser un avis ★",
    "Acheter / vendre (Pet's Shop)",
    "Modérer profils & annonces",
]
uc_x, uc_w, uc_h = 4.35, 4.6, 0.6
uc_y = [2.0 + i * 0.71 for i in range(6)]
for i, txt in enumerate(ucs):
    shape_text(s, MSO_SHAPE.OVAL, uc_x, uc_y[i], uc_w, uc_h, txt, WHITE, txt=DARK, size=10.5,
               bold=False, line=MED)
mid = [y + uc_h / 2 for y in uc_y]
# Associations
connect(s, 2.4, 2.46, uc_x, mid[0])                       # Visiteur → s'inscrire
for i in (1, 2, 3, 4):                                    # Membre → 4 cas
    connect(s, 2.4, 4.31, uc_x, mid[i])
connect(s, 10.75, 5.21, uc_x + uc_w, mid[5])              # Admin → modérer
para(tb(s, 2.95, 6.28, 7.4, 0.3), "L'administrateur est aussi un membre (il peut tout faire, plus la modération).",
     size=10.5, color=MED, italic=True, first=True)
notes(s, """
GLOSSAIRE :
• UML : « langage » de schémas standard pour décrire un logiciel de façon visuelle.
• Diagramme de cas d'utilisation : schéma UML qui montre les ACTEURS et les ACTIONS possibles sur le système.
• Acteur : un rôle d'utilisateur (ici Visiteur, Membre, Administrateur).
• Cas d'utilisation : une action que l'acteur peut réaliser (s'inscrire, publier, modérer…).
• Le grand cadre représente le système (le site) ; chaque trait relie un acteur à ce qu'il a le droit de faire.
• Modérer : vérifier et valider/supprimer des contenus (ici, l'admin approuve les annonces et surveille les profils).
""")

# ===========================================================================
# 9 · ÉCO-CONCEPTION
# ===========================================================================
s = content("Éco-conception", "Un site plus léger, plus rapide, plus sobre")
intro = tb(s, 0.7, 1.45, EMW - 1.4, 0.7)
para(intro, "Concevoir le site pour qu'il consomme moins de ressources (données, énergie) "
            "tout en restant rapide et accessible.", size=14, color=MED, italic=True, first=True)
colL = tb(s, 0.7, 2.35, 6.0, 4.0)
check(colL, "Images optimisées + chargement différé (lazy-load)", "elles ne se chargent que lorsqu'elles deviennent visibles.")
check(colL, "Aucune librairie lourde (ni Bootstrap, ni jQuery)", "moins de code à télécharger pour l'utilisateur.")
check(colL, "CSS factorisé avec SCSS", "variables et mixins évitent de répéter le même style.")
colR = tb(s, 7.0, 2.35, 5.6, 4.0)
check(colR, "Polices et icônes ciblées", "on n'importe que ce qui est réellement utilisé.")
check(colR, "Données chargées à la demande", "le serveur n'envoie que le nécessaire, au bon moment.")
check(colR, "Code structuré et réutilisable", "plus simple à maintenir et à faire évoluer.")
para(tb(s, 0.7, 6.45, EMW - 1.4, 0.5, anchor=MSO_ANCHOR.MIDDLE),
     "Bénéfices : pages plus rapides, moins de données mobiles consommées, meilleure accessibilité.",
     size=12.5, color=DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
notes(s, """
GLOSSAIRE :
• Éco-conception : démarche qui vise à réduire l'impact d'un site (poids des pages, énergie, données échangées).
• Lazy-load (chargement différé) : on ne télécharge une image que lorsqu'elle arrive à l'écran, pas avant.
• Librairie : ensemble de code tout prêt qu'on ajoute au projet ; certaines (jQuery, Bootstrap) alourdissent la page.
• Factoriser : écrire une seule fois un élément réutilisé partout, au lieu de le répéter.
• Mixin (SCSS) : bloc de styles réutilisable qu'on appelle à plusieurs endroits.
""")

# ===========================================================================
# 10 · SCSS (compétences)
# ===========================================================================
s = content("Le style avec SCSS", "Structurer et adapter l'affichage")
tf = tb(s, 0.7, 1.5, EMW - 1.4, 4.7)
check(tf, "Préprocesseur SCSS, architecture « 7-1 »", "le CSS est rangé en dossiers clairs, facile à maintenir.")
check(tf, "Variables et mixins", "couleurs et espacements cohérents partout, sans répétition.")
check(tf, "Mobile-first avec media queries", "l'affichage est pensé d'abord pour le téléphone, puis agrandi pour les écrans plus larges.")
check(tf, "Nommage BEM des classes", "des noms de classes lisibles et prévisibles.")
check(tf, "Modules modernes @use / @forward", "organisation propre, sans l'ancien @import déprécié.")
rncp_chip(s, "structurer le code avec un préprocesseur CSS  ·  agencement mobile-first (media queries)")
notes(s, """
GLOSSAIRE :
• CSS : langage qui décrit l'apparence d'une page (couleurs, tailles, positions).
• SCSS / préprocesseur : une version enrichie du CSS (variables, mixins, fichiers séparés) qui se « compile » ensuite en CSS classique.
• Architecture 7-1 : façon standard d'organiser les fichiers SCSS en 7 dossiers (base, composants, pages…) + 1 fichier principal.
• Variable : une valeur nommée réutilisable (ex. la couleur caramel) modifiable à un seul endroit.
• Mobile-first : on conçoit d'abord pour mobile, puis on ajoute des règles pour les grands écrans.
• Media query : règle CSS qui s'applique selon la taille de l'écran (ex. à partir de 768 px).
• BEM : convention de nommage des classes (Bloc__Élément--Modificateur) pour s'y retrouver.
• @use / @forward : instructions SCSS modernes pour relier les fichiers entre eux.
""")

# ===========================================================================
# 11 · HTML (compétences)
# ===========================================================================
s = content("La structure avec HTML", "Découper la maquette et intégrer les contenus")
tf = tb(s, 0.7, 1.5, EMW - 1.4, 4.7)
check(tf, "Découpage en blocs sémantiques", "header, nav, main, section, article, footer : une structure qui a du sens.")
check(tf, "Balises génériques (div, span) au bon endroit", "complètent le sémantique quand c'est pertinent.")
check(tf, "Accessibilité soignée", "lien d'évitement, aria-label, textes alternatifs, focus visible (objectif WCAG AA, validation W3C).")
check(tf, "Actifs intégrés", "images optimisées, vidéo et audio de présentation, polices importées.")
rncp_chip(s, "analyser & structurer la maquette en HTML  ·  typographie/accessibilité  ·  intégrer les actifs")
notes(s, """
GLOSSAIRE :
• HTML : langage qui structure le contenu d'une page (titres, paragraphes, images…).
• Balise sémantique : balise dont le nom décrit son rôle (header = en-tête, nav = navigation, article = contenu autonome). Cela aide moteurs de recherche et lecteurs d'écran.
• Balise générique : div et span, sans signification particulière, pour regrouper ou styliser.
• Accessibilité : rendre le site utilisable par tous, y compris les personnes en situation de handicap.
• Lien d'évitement (skip-link) : lien caché permettant d'aller directement au contenu principal au clavier.
• aria-label : texte invisible à l'écran mais lu par les lecteurs d'écran pour décrire un élément.
• Texte alternatif (alt) : description d'une image, lue si l'image ne s'affiche pas ou par un lecteur d'écran.
• WCAG AA : niveau de référence des normes d'accessibilité web.
• W3C : organisme qui définit les standards du web ; son validateur vérifie que le code HTML est correct.
• Actifs : ressources intégrées à la page (images, vidéos, sons, polices, scripts).
""")

# ===========================================================================
# 12 · JAVASCRIPT (compétences)
# ===========================================================================
s = content("L'interactivité avec JavaScript", "Rendre la page vivante et communiquer avec le serveur")
tf = tb(s, 0.7, 1.5, EMW - 1.4, 4.7)
check(tf, "Manipulation du DOM & gestion d'événements", "réagir aux clics, au clavier, à l'envoi de formulaires.")
check(tf, "Composants interactifs", "fenêtre modale, sélecteur d'étoiles, filtres, carrousel, réactions emoji.")
check(tf, "Validation des formulaires côté navigateur", "retours immédiats à l'utilisateur (ex. « cochez une note »).")
check(tf, "Requêtes HTTP asynchrones (fetch + async/await)", "le contenu se met à jour sans recharger la page.")
check(tf, "Dialogue avec mon back-end Express ET une API externe", "les deux cas prévus par la grille sont couverts.")
rncp_chip(s, "manipuler la page pour interagir avec l'utilisateur  ·  communication asynchrone avec un serveur")
notes(s, """
GLOSSAIRE :
• JavaScript : langage qui rend la page interactive (réagir aux actions de l'utilisateur).
• DOM : représentation, en mémoire, de la page HTML ; le JavaScript la lit et la modifie pour changer l'affichage.
• Événement : une action détectée (clic, touche clavier, envoi de formulaire) à laquelle on réagit.
• Modale : fenêtre qui s'ouvre par-dessus la page (ici, le détail d'un évènement).
• Requête HTTP : message envoyé au serveur pour demander ou envoyer des données.
• Asynchrone : l'action se fait en arrière-plan, sans bloquer ni recharger la page.
• fetch / async / await : les outils JavaScript modernes pour faire ces requêtes asynchrones.
• API externe : service tiers qui fournit des données (ici, une API d'animaux a servi à alimenter le site).
""")

# ===========================================================================
# 13 · QUALITÉ & SÉCURITÉ
# ===========================================================================
s = content("Qualité & sécurité — les points forts", "Un code soigné et protégé")
colL = tb(s, 0.7, 1.5, 6.0, 5.0)
para(colL, "Sécurité", size=16, color=BRAND, bold=True, font=HEAD, first=True, after=6)
bullet(colL, "jamais stockés en clair : « hachés » avec bcrypt + un secret (poivre).", lead="Mots de passe : ", size=12.5, after=7)
bullet(colL, "un jeton signé (JWT) prouve l'identité à chaque requête.", lead="Connexion : ", size=12.5, after=7)
bullet(colL, "requêtes SQL préparées : empêchent l'injection de code dans la base.", lead="Base de données : ", size=12.5, after=7)
bullet(colL, "tout texte affiché est échappé pour empêcher l'injection de scripts (XSS).", lead="Contenus : ", size=12.5, after=7)
bullet(colL, "en-têtes de sécurité (helmet) + limitation des tentatives (anti brute-force).", lead="Serveur : ", size=12.5, after=0)
colR = tb(s, 7.0, 1.5, 5.6, 5.0)
para(colR, "Qualité du code", size=16, color=BRAND, bold=True, font=HEAD, first=True, after=6)
bullet(colR, "rangé par domaine, commenté et indenté.", lead="Organisation : ", size=12.5, after=7)
bullet(colR, "noms explicites pour les variables, fichiers et URL.", lead="Lisibilité : ", size=12.5, after=7)
bullet(colR, "gestion centralisée des erreurs : pas de page d'erreur 500.", lead="Robustesse : ", size=12.5, after=7)
bullet(colR, "responsive : s'adapte du mobile au grand écran.", lead="Affichage : ", size=12.5, after=7)
bullet(colR, "espace d'administration avec un CRUD complet.", lead="Back-office : ", size=12.5, after=0)
notes(s, """
GLOSSAIRE :
• Hacher (hachage) : transformer un mot de passe en empreinte illisible et non réversible ; on ne stocke jamais le vrai mot de passe.
• bcrypt : algorithme de hachage spécialement conçu pour les mots de passe (volontairement lent, donc difficile à casser).
• Poivre (pepper) : secret ajouté avant le hachage, gardé hors de la base, pour renforcer la sécurité.
• JWT (jeton) : petit jeton signé envoyé à chaque requête pour prouver qu'on est bien connecté.
• Injection SQL : attaque qui glisse du code malveillant dans une requête à la base ; les « requêtes préparées » l'empêchent.
• Requête préparée : requête où les valeurs saisies sont traitées comme des données, jamais comme du code.
• XSS (cross-site scripting) : attaque qui injecte un script dans une page ; « échapper » le texte (neutraliser les caractères spéciaux) l'empêche.
• helmet : outil qui ajoute des en-têtes de sécurité au serveur.
• Brute-force : essayer des milliers de mots de passe ; la limitation de tentatives (rate-limiting) freine cette attaque.
• Erreur 500 : erreur serveur ; une bonne gestion des erreurs évite que le site « plante ».
• CRUD : les 4 opérations de base sur des données — Créer, Lire, Modifier, Supprimer (Create, Read, Update, Delete).
• Responsive : qui s'adapte automatiquement à la taille de l'écran.
• Back-office : l'espace d'administration, réservé aux gestionnaires du site.
""")

# ===========================================================================
# 14 · REMERCIEMENTS + DÉMO
# ===========================================================================
s = slide()
rect(s, 0, 0, EMW, EMH, DARK)
rect(s, 0, 0, EMW, 0.35, BRAND); rect(s, 0, EMH - 0.35, EMW, 0.35, BRAND)
if os.path.exists(LOGO):
    img_fit(s, LOGO, EMW / 2 - 0.95, 1.1, 1.9, 1.6, border=None)
t = tb(s, 1.2, 3.0, EMW - 2.4, 2.2)
para(t, "Merci de votre attention", size=40, color=WHITE, bold=True, font=HEAD,
     align=PP_ALIGN.CENTER, first=True, after=10)
para(t, "Place à la démonstration en direct", size=20, color=BRAND, bold=True, align=PP_ALIGN.CENTER, after=14)
para(t, "Je vous propose de parcourir ensemble le site : création de profil, feed, "
        "évènements et avis, marketplace, et espace administrateur.",
     size=14, color=CREAM, align=PP_ALIGN.CENTER, after=0)
notes(s, """
Diapo de clôture — enchaîner sur la démonstration (≈ 7 min) : création d'un profil animal, feed,
évènements (inscription + avis), Pet's Shop, et espace administrateur (CRUD).
Pas de terme technique nouveau ici.
""")

# --- Sauvegarde -------------------------------------------------------------
OUT = os.path.join(ROOT, "presentation", "Pet-s-Book-Soutenance-Bloc2.pptx")
try:
    prs.save(OUT)
    saved = OUT
except PermissionError:
    saved = OUT.replace(".pptx", "-v2.pptx")
    prs.save(saved)
print(f"OK - {len(prs.slides)} diapositives")
print(f"Fichier : {saved}")
